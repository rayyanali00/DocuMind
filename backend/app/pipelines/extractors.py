"""Text extraction for supported formats.

Returns an ExtractionResult that flags whether OCR was used so the caller
(ingest task) can mark the File row accordingly (FR-025).
"""
from __future__ import annotations

import csv
from dataclasses import dataclass
from pathlib import Path

from loguru import logger

from app.pipelines.ocr import ocr_image, ocr_pdf, pdf_appears_scanned
from app.pipelines.validators import IMAGE_EXTENSIONS, get_extension


class ExtractionError(Exception):
    pass


@dataclass(slots=True)
class ExtractionResult:
    text: str
    ocr_used: bool = False


def _extract_txt(path: Path) -> ExtractionResult:
    return ExtractionResult(text=path.read_text(encoding="utf-8", errors="replace"))


def _extract_pdf(path: Path) -> ExtractionResult:
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
        parts: list[str] = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception as exc:  # noqa: BLE001
                logger.warning(f"PDF page extraction failed for {path}: {exc}")
        text = "\n".join(parts).strip()
    except Exception as exc:
        logger.warning(f"PDF text extraction failed entirely, falling back to OCR: {exc}")
        text = ""

    if pdf_appears_scanned(text):
        logger.info(f"PDF appears scanned, running OCR: {path.name}")
        ocr_text = ocr_pdf(path)
        return ExtractionResult(text=ocr_text, ocr_used=True)
    return ExtractionResult(text=text)


def _extract_docx(path: Path) -> ExtractionResult:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return ExtractionResult(text="\n".join(parts).strip())


def _extract_csv(path: Path) -> ExtractionResult:
    rows: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        for row in csv.reader(f):
            rows.append(", ".join(row))
    return ExtractionResult(text="\n".join(rows).strip())


def _extract_xlsx(path: Path) -> ExtractionResult:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(", ".join(cells))
    wb.close()
    return ExtractionResult(text="\n".join(parts).strip())


def _extract_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return ExtractionResult(text="\n".join(parts).strip())


def _extract_image(path: Path) -> ExtractionResult:
    return ExtractionResult(text=ocr_image(path), ocr_used=True)


_EXTRACTORS = {
    "txt": _extract_txt,
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "csv": _extract_csv,
    "xlsx": _extract_xlsx,
    "pptx": _extract_pptx,
}


def extract_text(path: Path) -> ExtractionResult:
    ext = get_extension(path.name)
    if ext in IMAGE_EXTENSIONS:
        return _extract_image(path)
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ExtractionError(f"No extractor registered for .{ext}")
    try:
        return extractor(path)
    except Exception as exc:
        raise ExtractionError(f"Failed to extract text from {path.name}: {exc}") from exc
